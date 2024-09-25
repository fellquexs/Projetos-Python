from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import openpyxl
import os
from datetime import datetime

#Elaborado por @Gustavo Felix - 21.05.24

def carregar_dados_excel(resultados_AT):
    wb = openpyxl.load_workbook(resultados_AT)
    planilha = wb.active
    custo_gerador = str(planilha["T150"].value)
    periodo_analise = str(planilha["T139"].value)
    data_estudo = str(planilha["T137"].value)
    remuneracao = str(planilha["T152"].value)
    nome_uc = f'Unidade {str(planilha["S230"].value)}'
    EC_AT_VD = f'R${"{:,.2f}".format(float(planilha["T317"].value)).replace(",", ".")}'
    EC_AT_MD = f'R${"{:,.2f}".format(float(planilha["T327"].value)).replace(",", ".")}'
    PERCENT_TXT_EC_AT_VD = f'{float(planilha["T318"].value)*100:.0f}%'
    PERCENT_TXT_EC_AT_MD = f'{float(planilha["T328"].value)*100:.0f}%'
       
    # Verifica se o valor de CNPJ está presente e aplica a formatação
    cnpj = str(planilha["T230"].value)
        # Formata o CNPJ (##.###.###/####-##)
    cnpj_formatado = f"{cnpj[:2]}.{cnpj[2:5]}.{cnpj[5:8]}/{cnpj[8:12]}-{cnpj[12:]}"



    dados_tabela_tarifas = [
        ["Tarifa Cativa de Energia Vigente", "TE Ponta (R$/MWh)", "TE Fora Ponta (R$/MWh)"],
        [str(planilha["S194"].value), f'R${str(planilha["T194"].value).replace(".", ",")}', f'R${str(planilha["U194"].value).replace(".", ",")}']
    ]

    dados_tabela_imp = [
        ["ICMS", "Alíquota", "Base de Cálculo na TUSD"],
        [str(planilha["W194"].value), f'{float(planilha["X194"].value) * 100:.0f}%', str(planilha["Y194"].value)]
    ]
    dados_tabela_preco_equilibrio_AT = [
        ["Preço de Equilíbrio", str(planilha["T182"].value), str(planilha["U182"].value), str(planilha["V182"].value), str(planilha["W182"].value), str(planilha["X182"].value)],
        ["Preço R$/MWh - Atacadista", f'R${float(planilha["T183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["U183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["V183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["W183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["X183"].value):,.2f}'.replace(".", ",")]
    ]

    dados_tabela_preco_livre_AT = [
        ["Preço Mercado Livre", str(planilha["T185"].value), str(planilha["U185"].value), str(planilha["V185"].value), str(planilha["W185"].value), str(planilha["X185"].value)],
        ["Energia i50% R$/MWh - Atacadista", f'R${float(planilha["T186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["U186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["V186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["W186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["X186"].value):,.2f}'.replace(".", ",")]
    ]

    dados_tabela_reajustes = [
        ["Índice de Reajuste", str(planilha["T188"].value), str(planilha["U188"].value), str(planilha["V188"].value), str(planilha["W188"].value), str(planilha["X188"].value)],
        [str(planilha["S189"].value), f'{float(planilha["T189"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["U189"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["V189"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["W189"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["X189"].value)*100:.2f}%'.replace(".", ",")],
        [str(planilha["S190"].value), f'{float(planilha["T190"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["U190"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["V190"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["W190"].value)*100:.2f}%'.replace(".", ","), f'{float(planilha["X190"].value)*100:.2f}%'.replace(".", ",")]
    ]

    dados_tabela_info_uc = [
        ["Unidade", "CNPJ", "Distribuidora", "Submercado", "Mod. Tarifária", "Energia Contratada", "Demanda Cativo (KW)", "Demanda Livre (KW)","Demanda Livre Varejista (KW)", "Inicio de Fornecimento", "% Consumo Ponta", "Volume Médio de Consumo (MWm)"],
        [str(planilha["S230"].value), cnpj_formatado, str(planilha["U230"].value), str(planilha["V230"].value), str(planilha["W230"].value), str(planilha["X230"].value), str(planilha["Y230"].value), str(planilha["Z230"].value), str(planilha["Y230"].value) , str(planilha["AA230"].value), f'{float(planilha["AO230"].value)*100:,.2f}%'.replace(".", ","), f'{float(planilha["AB230"].value):,.3f}'.replace(".", ",")]
    ]

    dados_tabela_AT_VED = [
    ["PERÍODO", "CUSTO CATIVO MENSAL", "CUSTO LIVRE MENSAL", "ECONOMIA MENSAL", "ECONOMIA ANUAL", "ECONOMIA (%)"],
    [str(planilha["T446"].value), f'R${"{:,.2f}".format(float(planilha["U446"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V446"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W446"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X446"].value)).replace(",", ".")}', f'{float(planilha["Y446"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T447"].value), f'R${"{:,.2f}".format(float(planilha["U447"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V447"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W447"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X447"].value)).replace(",", ".")}', f'{float(planilha["Y447"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T448"].value), f'R${"{:,.2f}".format(float(planilha["U448"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V448"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W448"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X448"].value)).replace(",", ".")}', f'{float(planilha["Y448"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T449"].value), f'R${"{:,.2f}".format(float(planilha["U449"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V449"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W449"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X449"].value)).replace(",", ".")}', f'{float(planilha["Y449"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T450"].value), f'R${"{:,.2f}".format(float(planilha["U450"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V450"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W450"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X450"].value)).replace(",", ".")}', f'{float(planilha["Y450"].value)*100:.0f}%'.replace(".", ",")]
    ]


    dados_tabela_AT_MED = [
        ["PERÍODO", "CUSTO CATIVO MENSAL", "CUSTO LIVRE MENSAL", "ECONOMIA MENSAL", "ECONOMIA ANUAL", "ECONOMIA (%)"],
        [str(planilha["T461"].value), f'R${"{:,.2f}".format(float(planilha["U461"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V461"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W461"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X461"].value)).replace(",", ".")}', f'{float(planilha["Y461"].value)*100:.0f}%'],
        [str(planilha["T462"].value), f'R${"{:,.2f}".format(float(planilha["U462"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V462"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W462"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X462"].value)).replace(",", ".")}', f'{float(planilha["Y462"].value)*100:.0f}%'],
        [str(planilha["T463"].value), f'R${"{:,.2f}".format(float(planilha["U463"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V463"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W463"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X463"].value)).replace(",", ".")}', f'{float(planilha["Y463"].value)*100:.0f}%'],
        [str(planilha["T464"].value), f'R${"{:,.2f}".format(float(planilha["U464"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V464"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W464"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X464"].value)).replace(",", ".")}', f'{float(planilha["Y464"].value)*100:.0f}%'],
        [str(planilha["T465"].value), f'R${"{:,.2f}".format(float(planilha["U465"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V465"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W465"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X465"].value)).replace(",", ".")}', f'{float(planilha["Y465"].value)*100:.0f}%']
    ]

    if isinstance(data_estudo, datetime):
        data_estudo = data_estudo.strftime("%d-%m-%y")
    return custo_gerador, periodo_analise, data_estudo, remuneracao, nome_uc , dados_tabela_tarifas, dados_tabela_imp, dados_tabela_preco_equilibrio_AT, dados_tabela_preco_livre_AT, dados_tabela_reajustes, dados_tabela_info_uc, dados_tabela_AT_VED, dados_tabela_AT_MED, EC_AT_VD, EC_AT_MD, PERCENT_TXT_EC_AT_VD, PERCENT_TXT_EC_AT_MD

def preencher_substituicoes(prs, substituicoes, cor_montserrat):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in substituicoes:
                shape.text_frame.text = substituicoes[shape.name]
                for paragraph in shape.text_frame.paragraphs:
                    # Condição para centralizar texto
                    if shape.name in ["EC_AT_VD", "EC_AT_MD"]:
                        paragraph.alignment = PP_ALIGN.CENTER  # Centraliza o parágrafo
                    # Condição para centralizar e alterar formatação específica
                    if shape.name in ["PERCENT_EC_AT_VD", "PERCENT_EC_AT_MD"]:
                        paragraph.alignment = PP_ALIGN.CENTER  # Centraliza o parágrafo
                        for run in paragraph.runs:
                            run.font.size = Pt(24)  # Define o tamanho da fonte para 24
                            run.font.color.rgb = RGBColor(89, 89, 90)  # Define a cor da fonte para RGB(89, 89, 90)
                            run.font.bold = True
                    else:
                        for run in paragraph.runs:
                            run.font.name = "Montserrat"
                            run.font.size = Pt(18)
                            run.font.color.rgb = cor_montserrat
                            run.font.bold = True
    # Tratamento especial para a variável nome_uc
    if "nome_uc" in substituicoes:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name == "nome_uc":
                    # Separar "Unidade [Nome cliente]" em partes
                    texto = substituicoes["nome_uc"]
                    partes = texto.split(' ')
                    unidade = partes[0].strip()
                    nome_cliente = partes[1].strip(']')

                    # Definir o texto no shape
                    shape.text_frame.clear()  # Limpar qualquer texto existente
                    p = shape.text_frame.add_paragraph()
                    
                    # Adicionar "Unidade" ao parágrafo
                    run_unidade = p.add_run()
                    run_unidade.text = unidade + ' '
                    run_unidade.font.name = "Montserrat"
                    run_unidade.font.size = Pt(18)
                    run_unidade.font.color.rgb = RGBColor(0, 0, 0)  # preto
                    run_unidade.font.bold = False

                    # Adicionar "Nome cliente" ao parágrafo
                    run_nome_cliente = p.add_run()
                    run_nome_cliente.text = nome_cliente
                    run_nome_cliente.font.name = "Montserrat"
                    run_nome_cliente.font.size = Pt(18)
                    run_nome_cliente.font.color.rgb = RGBColor(0, 0, 0)  # preto
                    run_nome_cliente.font.bold = True

def preencher_tabela(tabela, dados, formatar=False):
    for i, linha in enumerate(dados):
        for j, valor in enumerate(linha):
            cell = tabela.cell(i, j)
            cell.text = valor
            # Ajusta o tamanho da fonte e alinha o texto
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    # Aplica a formatação em negrito se for a coluna "PERÍODO" e não for a linha do cabeçalho
                    if formatar and j == 0 and i > 0:
                        run.font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def preencher_tabelas(prs, dados_tabelas,tabelas_para_formatar):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name in dados_tabelas:
                tabela = shape.table
                dados = dados_tabelas[shape.name]
                # Verifica se a tabela deve ser formatada
                formatar = shape.name in tabelas_para_formatar                
                preencher_tabela(tabela, dados, formatar)


path = os.path.dirname(os.path.abspath(__file__))
nome_arquivo_pptx = os.path.join(path, "AT x VJPF - 1 UNIDADE.pptx")
resultados_AT = os.path.join(path, "AT.xlsx")

custo_gerador, periodo_analise, data_estudo, remuneracao, nome_uc , dados_tabela_tarifas, dados_tabela_imp, dados_tabela_preco_equilibrio_AT, dados_tabela_preco_livre_AT, dados_tabela_reajustes, dados_tabela_info_uc, dados_tabela_AT_VED, dados_tabela_AT_MED, EC_AT_VD, EC_AT_MD, PERCENT_TXT_EC_AT_VD, PERCENT_TXT_EC_AT_MD = carregar_dados_excel(resultados_AT)

substituicoes = {
    "Retângulo 4": f'{data_estudo}',
    "Retângulo 5": f'{periodo_analise}',
    "Retângulo 43": f'{custo_gerador}',
    "Retângulo 139": f'{remuneracao}',
    "nome_uc": f'{nome_uc}',
    "EC_AT_VD": f'{EC_AT_VD}',
    "EC_AT_MD": f'{EC_AT_MD}',
    "PERCENT_EC_AT_VD": f'{PERCENT_TXT_EC_AT_VD}',
    "PERCENT_EC_AT_MD": f'{PERCENT_TXT_EC_AT_MD}'
}

dados_tabelas = {
    "TabTarifas": dados_tabela_tarifas,
    "TabImp": dados_tabela_imp,
    "TabPrecoeqAT": dados_tabela_preco_equilibrio_AT,
    "TabPrecomlAT": dados_tabela_preco_livre_AT,
    "TabReaj": dados_tabela_reajustes,
    "Info_UC": dados_tabela_info_uc,
    "Tab_ATVD": dados_tabela_AT_VED,
    "Tab_ATMD": dados_tabela_AT_MED
}

#SEGUNDO ARQUIVO VJDG
def carregar_dados_excel_segundo_arquivo(resultados_VJPF):
    wb = openpyxl.load_workbook(resultados_VJPF)
    planilha = wb.active

    EC_VJPF_VD = f'R${"{:,.2f}".format(float(planilha["T317"].value)).replace(",", ".")}'
    EC_VJPF_MD = f'R${"{:,.2f}".format(float(planilha["T327"].value)).replace(",", ".")}'
    PERCENT_TXT_EC_VJPF_VD = f'{float(planilha["T318"].value)*100:.0f}%'
    PERCENT_TXT_EC_VJPF_MD = f'{float(planilha["T328"].value)*100:.0f}%'

    
    dados_tabela_preco_equilibrio_VJPF = [
        ["Preço R$/MWh - Varejista", f'R${float(planilha["T183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["U183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["V183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["W183"].value):,.2f}'.replace(".", ","), f'R${float(planilha["X183"].value):,.2f}'.replace(".", ",")]
    ]

    dados_tabela_preco_livre_VJ = [
        ["Energia i50% R$/MWh - Varejista", f'R${float(planilha["T186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["U186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["V186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["W186"].value):,.2f}'.replace(".", ","), f'R${float(planilha["X186"].value):,.2f}'.replace(".", ",")]
    ]


    dados_tabela_VJPF_VED = [
    ["PERÍODO", "CUSTO CATIVO MENSAL", "CUSTO LIVRE MENSAL", "ECONOMIA MENSAL", "ECONOMIA ANUAL", "ECONOMIA (%)"],
    [str(planilha["T446"].value), f'R${"{:,.2f}".format(float(planilha["U446"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V446"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W446"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X446"].value)).replace(",", ".")}', f'{float(planilha["Y446"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T447"].value), f'R${"{:,.2f}".format(float(planilha["U447"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V447"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W447"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X447"].value)).replace(",", ".")}', f'{float(planilha["Y447"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T448"].value), f'R${"{:,.2f}".format(float(planilha["U448"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V448"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W448"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X448"].value)).replace(",", ".")}', f'{float(planilha["Y448"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T449"].value), f'R${"{:,.2f}".format(float(planilha["U449"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V449"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W449"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X449"].value)).replace(",", ".")}', f'{float(planilha["Y449"].value)*100:.0f}%'.replace(".", ",")],
    [str(planilha["T450"].value), f'R${"{:,.2f}".format(float(planilha["U450"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V450"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W450"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X450"].value)).replace(",", ".")}', f'{float(planilha["Y450"].value)*100:.0f}%'.replace(".", ",")]
    ]


    dados_tabela_VJPF_MED = [
        ["PERÍODO", "CUSTO CATIVO MENSAL", "CUSTO LIVRE MENSAL", "ECONOMIA MENSAL", "ECONOMIA ANUAL", "ECONOMIA (%)"],
        [str(planilha["T461"].value), f'R${"{:,.2f}".format(float(planilha["U461"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V461"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W461"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X461"].value)).replace(",", ".")}', f'{float(planilha["Y461"].value)*100:.0f}%'],
        [str(planilha["T462"].value), f'R${"{:,.2f}".format(float(planilha["U462"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V462"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W462"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X462"].value)).replace(",", ".")}', f'{float(planilha["Y462"].value)*100:.0f}%'],
        [str(planilha["T463"].value), f'R${"{:,.2f}".format(float(planilha["U463"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V463"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W463"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X463"].value)).replace(",", ".")}', f'{float(planilha["Y463"].value)*100:.0f}%'],
        [str(planilha["T464"].value), f'R${"{:,.2f}".format(float(planilha["U464"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V464"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W464"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X464"].value)).replace(",", ".")}', f'{float(planilha["Y464"].value)*100:.0f}%'],
        [str(planilha["T465"].value), f'R${"{:,.2f}".format(float(planilha["U465"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["V465"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["W465"].value)).replace(",", ".")}', f'R${"{:,.2f}".format(float(planilha["X465"].value)).replace(",", ".")}', f'{float(planilha["Y465"].value)*100:.0f}%']
    ]
    return dados_tabela_VJPF_VED, dados_tabela_VJPF_MED, EC_VJPF_VD, EC_VJPF_MD, PERCENT_TXT_EC_VJPF_VD, PERCENT_TXT_EC_VJPF_MD, dados_tabela_preco_equilibrio_VJPF, dados_tabela_preco_livre_VJ

resultados_VJPF = os.path.join(path, "VJPF.xlsx")


def preencher_substituicoes_segundo_arquivo(prs, substituicoes_segundo_arquivo, cor_montserrat):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name in substituicoes_segundo_arquivo:
                shape.text_frame.text = substituicoes_segundo_arquivo[shape.name]
                for paragraph in shape.text_frame.paragraphs:
                    # Condição para centralizar texto
                    if shape.name in ["EC_VJPF_VD", "EC_VJPF_MD"]:
                        paragraph.alignment = PP_ALIGN.CENTER  # Centraliza o parágrafo
                    # Condição para centralizar e alterar formatação específica
                    if shape.name in ["PERCENT_EC_VJPF_VD", "PERCENT_EC_VJPF_MD"]:
                        paragraph.alignment = PP_ALIGN.CENTER  # Centraliza o parágrafo
                        for run in paragraph.runs:
                            run.font.size = Pt(24)  # Define o tamanho da fonte para 24
                            run.font.color.rgb = RGBColor(89, 89, 90)  # Define a cor da fonte para RGB(89, 89, 90)
                            run.font.bold = True
                    else:
                        for run in paragraph.runs:
                            run.font.name = "Montserrat"
                            run.font.size = Pt(18)
                            run.font.color.rgb = cor_montserrat
                            run.font.bold = True




# Função para preencher a tabela com formatação específica
def preencher_tabela_segundo_arquivo(tabela, dados, formatar=False, tabela_nome=None):
    for i, linha in enumerate(dados):
        for j, valor in enumerate(linha):
            cell = tabela.cell(i, j)
            cell.text = valor
            # Ajusta o tamanho da fonte e alinha o texto
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if tabela_nome in ["TabPrecoeqVJ", "TabPrecomlVJ"]:
                        run.font.name = "Montserrat"
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Define a cor da fonte para preto
                        run.font.bold = False  # Remove o negrito
                    else:
                        run.font.size = Pt(10)
                        # Aplica a formatação em negrito se for a coluna "PERÍODO" e não for a linha do cabeçalho
                        if formatar and j == 0 and i > 0:
                            run.font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Função para preencher as tabelas na apresentação
def preencher_tabelas_segundo_arquivo(prs, dados_tabelas_segundo_arquivo, tabelas_para_formatar):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table and shape.name in dados_tabelas_segundo_arquivo:
                tabela = shape.table
                dados = dados_tabelas_segundo_arquivo[shape.name]
                # Verifica se a tabela deve ser formatada
                formatar = shape.name in tabelas_para_formatar
                preencher_tabela_segundo_arquivo(tabela, dados, formatar, tabela_nome=shape.name)


dados_tabela_VJPF_VED, dados_tabela_VJPF_MED,EC_VJPF_VD, EC_VJPF_MD, PERCENT_TXT_EC_VJPF_VD, PERCENT_TXT_EC_VJPF_MD, dados_tabela_preco_equilibrio_VJPF, dados_tabela_preco_livre_VJ = carregar_dados_excel_segundo_arquivo(resultados_VJPF)





dados_tabelas_segundo_arquivo = {
    "Tab_VJPFVD": dados_tabela_VJPF_VED,
    "Tab_VJPFMD": dados_tabela_VJPF_MED,
    "Tab_ATVD": dados_tabela_AT_VED, #essas duas
    "Tab_ATMD": dados_tabela_AT_MED, #essa tmb
    "TabPrecoeqVJ": dados_tabela_preco_equilibrio_VJPF,
    "TabPrecomlVJ":dados_tabela_preco_livre_VJ
}


# Lista de nomes das tabelas que devem ser formatadas (use os mesmos nomes que no dicionário)
tabelas_para_formatar = [
    "Tab_VJPFVD",
    "Tab_VJPFMD",
    "Tab_ATVD",
    "Tab_ATMD"
    # Adicione os nomes das outras duas tabelas aqui
]



substituicoes_segundo_arquivo = {
    "EC_VJPF_VD": f'{EC_VJPF_VD}',
    "EC_VJPF_MD": f'{EC_VJPF_MD}',
    "PERCENT_EC_VJPF_VD": f'{PERCENT_TXT_EC_VJPF_VD}',
    "PERCENT_EC_VJPF_MD": f'{PERCENT_TXT_EC_VJPF_MD}'

}

cor_montserrat = RGBColor(255, 147, 0)  

prs = Presentation(nome_arquivo_pptx)

preencher_substituicoes(prs, substituicoes, cor_montserrat)
preencher_tabelas(prs, dados_tabelas,tabelas_para_formatar)
preencher_tabelas_segundo_arquivo(prs, dados_tabelas_segundo_arquivo, tabelas_para_formatar)
preencher_substituicoes_segundo_arquivo(prs,substituicoes_segundo_arquivo,cor_montserrat)


saida_name = f'{nome_uc[8:]} COMPARATIVO ATXVJPF.pptx'
nome_arquivo_saida = os.path.join(path, saida_name)

#nome_arquivo_saida = os.path.join(path, "COMPARATIVO_VJPFXVJDG.pptx")
prs.save(nome_arquivo_saida)
